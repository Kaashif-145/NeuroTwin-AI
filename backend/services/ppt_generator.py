import os
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _sanitize_text(text):
    if not isinstance(text, str):
        return ""
    
    # Mathematical character mapping for PPT
    replacements = {
        "∫": "INTEGRAL", "∑": "SUM", "≈": "≈", "≤": "≤", "≥": "≥",
        "±": "±", "→": "→", "∞": "∞", "π": "π", "θ": "θ",
        "α": "α", "β": "β", "γ": "γ", "δ": "δ", "Δ": "Δ",
        "≠": "≠", "λ": "λ", "μ": "μ", "σ": "σ", "ω": "ω", "×": "×", "÷": "÷", "√": "√",
        "²": "²", "³": "³", "¼": "¼", "½": "½", "¾": "¾"
    }
    
    for search, replace in replacements.items():
        text = text.replace(search, replace)

    text = re.sub(r"\s+", " ", text).strip()
    return text


def _summary_sections(summary):
    blocks = [block.strip() for block in re.split(r"\n\s*\n", summary or "") if block.strip()]
    sections = []
    for block in blocks:
        lines = [line.strip() for line in block.splitlines() if line.strip()]
        if not lines:
            continue
        title = lines[0].rstrip(":")
        body = " ".join(lines[1:]) if len(lines) > 1 else lines[0]
        sections.append((title, body))
    return sections or [("Detailed Summary", _sanitize_text(summary))]


def move_to_back(shape):
    parent = shape._element.getparent()
    parent.insert(0, shape._element)


def add_background(slide, prs, color=(244, 248, 255), accent=(213, 230, 255)):
    background = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(*color)
    background.line.fill.background()
    move_to_back(background)

    accent_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        prs.slide_width - Inches(2.5),
        Inches(-0.6),
        Inches(3.2),
        Inches(3.2),
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = RGBColor(*accent)
    accent_shape.line.fill.background()
    move_to_back(accent_shape)

    lower_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        prs.slide_height - Inches(1.15),
        prs.slide_width,
        Inches(1.15),
    )
    lower_band.fill.solid()
    lower_band.fill.fore_color.rgb = RGBColor(255, 214, 153)
    lower_band.line.fill.background()
    move_to_back(lower_band)

    corner_glow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(-0.7),
        prs.slide_height - Inches(2.3),
        Inches(3.2),
        Inches(3.2),
    )
    corner_glow.fill.solid()
    corner_glow.fill.fore_color.rgb = RGBColor(255, 179, 102)
    corner_glow.line.fill.background()
    move_to_back(corner_glow)


def add_text_box(
    slide,
    text,
    left,
    top,
    width,
    height,
    font_size=18,
    color=(33, 43, 54),
    bold=False,
    alignment=PP_ALIGN.LEFT,
):
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = _sanitize_text(text)
    p.font.size = Pt(font_size)
    p.font.name = 'Cambria'
    p.font.color.rgb = RGBColor(*color)
    p.font.bold = bold
    p.alignment = alignment
    return tx_box


def add_bullet_box(slide, title, items, left, top, width, height, fill=(255, 255, 255)):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(255, 196, 112)

    text_frame = shape.text_frame
    text_frame.clear()
    title_paragraph = text_frame.paragraphs[0]
    title_paragraph.text = _sanitize_text(title)
    title_paragraph.font.size = Pt(18)
    title_paragraph.font.name = 'Cambria'
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(18, 43, 85)

    for item in items:
        p = text_frame.add_paragraph()
        p.text = _sanitize_text(item)
        p.level = 0
        p.font.size = Pt(12)
        p.font.name = 'Cambria'
        p.font.color.rgb = RGBColor(55, 55, 55)
        p.bullet = True


def add_image_safely(slide, image_path, left, top, width=None, height=None):
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)


def generate_presentation(
    title,
    summary,
    topics,
    output_path="data/summary_presentation.pptx",
    glossary=None,
    equations=None,
    key_features=None,
):
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs = Presentation()

    clean_topics = [_sanitize_text(topic) for topic in topics[:10]]
    clean_equations = [_sanitize_text(eq) for eq in (equations or [])[:6]]
    clean_features = [_sanitize_text(feature) for feature in (key_features or [])[:6]]
    clean_glossary = glossary or {}
    sections = _summary_sections(summary)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs, color=(255, 247, 237), accent=(255, 204, 153))
    add_image_safely(slide, "frontend/assets/hero.png", Inches(5.45), Inches(0.45), width=Inches(3.55))
    add_text_box(slide, "Detailed Study Deck", Inches(0.6), Inches(0.95), Inches(4.5), Inches(0.8), 28, (77, 32, 12), True)
    add_text_box(slide, title, Inches(0.6), Inches(1.9), Inches(4.8), Inches(1.0), 22, (173, 77, 31), True)
    add_text_box(
        slide,
        "Visual summary slides covering major concepts, key features, detailed explanations, and mathematical equations.",
        Inches(0.6),
        Inches(3.0),
        Inches(4.7),
        Inches(1.5),
        15,
        (98, 71, 58),
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs, color=(255, 249, 240), accent=(255, 220, 163))
    add_text_box(slide, "Main Features", Inches(0.6), Inches(0.5), Inches(4), Inches(0.7), 24, (77, 32, 12), True)
    add_image_safely(slide, "frontend/assets/concept.png", Inches(6.95), Inches(0.55), width=Inches(2.05))
    add_bullet_box(
        slide,
        "Top Highlights",
        clean_features or clean_topics[:6] or ["Key features will appear here after processing."],
        Inches(0.6),
        Inches(1.3),
        Inches(8.2),
        Inches(5.5),
        fill=(255, 252, 247),
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs, color=(245, 250, 255), accent=(173, 220, 255))
    add_text_box(slide, "Core Concepts", Inches(0.6), Inches(0.5), Inches(4), Inches(0.7), 24, (14, 65, 103), True)
    left_topics = clean_topics[:5]
    right_topics = clean_topics[5:10]
    add_bullet_box(slide, "Concept Set A", left_topics or ["No topics detected"], Inches(0.6), Inches(1.4), Inches(4.0), Inches(4.8), fill=(255, 255, 255))
    add_bullet_box(slide, "Concept Set B", right_topics or ["Additional topics will appear here"], Inches(5.0), Inches(1.4), Inches(4.0), Inches(4.8), fill=(239, 248, 255))

    for index, (section_title, section_body) in enumerate(sections[:3], start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(
            slide,
            prs,
            color=(255, 248, 243) if index % 2 else (245, 251, 255),
            accent=(255, 205, 170) if index % 2 else (180, 226, 255),
        )
        add_text_box(slide, f"Detailed Summary {index}", Inches(0.6), Inches(0.45), Inches(4.5), Inches(0.7), 24, (77, 32, 12), True)
        add_bullet_box(
            slide,
            section_title,
            re.split(r"(?<=[.!?])\s+", section_body)[:6] or [section_body],
            Inches(0.6),
            Inches(1.3),
            Inches(8.2),
            Inches(5.4),
            fill=(255, 253, 250),
        )

    if clean_glossary:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide, prs, color=(244, 252, 248), accent=(173, 232, 202))
        add_text_box(slide, "Glossary", Inches(0.6), Inches(0.5), Inches(3.5), Inches(0.7), 24, (24, 83, 62), True)
        glossary_items = [f"{key}: {value}" for key, value in list(clean_glossary.items())[:6]]
        add_bullet_box(slide, "Major Terms", glossary_items, Inches(0.6), Inches(1.3), Inches(8.2), Inches(5.4), fill=(250, 255, 252))

    if clean_equations:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide, prs, color=(244, 246, 255), accent=(196, 201, 255))
        add_text_box(slide, "Mathematical Equations", Inches(0.6), Inches(0.5), Inches(5.2), Inches(0.7), 24, (42, 46, 120), True)
        add_bullet_box(
            slide,
            "Detected Expressions",
            clean_equations,
            Inches(0.6),
            Inches(1.3),
            Inches(8.2),
            Inches(5.4),
            fill=(250, 250, 255),
        )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs, color=(255, 244, 236), accent=(255, 194, 143))
    add_text_box(slide, "Ready for Revision", Inches(0.6), Inches(2.2), Inches(8.0), Inches(0.8), 28, (120, 49, 14), True, PP_ALIGN.CENTER)
    add_text_box(
        slide,
        "Use this deck with the PDF report for detailed review, concept recall, and equation practice.",
        Inches(1.0),
        Inches(3.2),
        Inches(7.2),
        Inches(1.2),
        16,
        (98, 71, 58),
        alignment=PP_ALIGN.CENTER,
    )

    prs.save(output_path)
    return output_path
